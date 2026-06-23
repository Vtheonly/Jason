import logging
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("computer-science-slides")

class ComputerScienceSlides:
    @staticmethod
    def draw_source_code_listing_box(slide, left_emu, top_emu, width_emu, height_emu, source_code_string):
        logger.info("Generating Computer Science source code listing block.")
        shapes = slide.shapes
        
        # 1. Draw a dark IDE background container shape
        ide_card = shapes.add_shape(MSO_SHAPE.RECTANGLE, left_emu, top_emu, width_emu, height_emu)
        ide_card.fill.solid()
        ide_card.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#1E1E1E")
        ide_card.line.color.rgb = ColorTranslator.hex_to_rgb("#333333")
        
        # 2. Add source code string to IDE code block shape
        tf = ide_card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_top = Pt(10)
        
        p = tf.paragraphs[0]
        p.text = source_code_string
        
        for run in p.runs:
            run.font.name = "Consolas"
            run.font.size = Pt(9.5)
            run.font.color.rgb = ColorTranslator.hex_to_rgb("#D4D4D4")
            
        logger.info("Source code listing block rendered successfully.")
        return ide_card