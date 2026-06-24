import logging
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("engineering-thesis-slides")

class EngineeringThesisSlides:
    @staticmethod
    def draw_constraints_matrix_block(slide, left_emu, top_emu, width_emu, height_emu, design_constraints_list):
        logger.info(f"Generating engineering thesis design constraints checklist card. Items: {len(design_constraints_list)}")
        shapes = slide.shapes
        
        # Draw background constraints matrix container card
        bg_card = shapes.add_shape(MSO_SHAPE.RECTANGLE, left_emu, top_emu, width_emu, height_emu)
        bg_card.fill.solid()
        bg_card.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#ECEFF1")
        bg_card.line.color.rgb = ColorTranslator.hex_to_rgb("#CFD8DC")
        bg_card.line.width = Pt(1.5)
        
        tf = bg_card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(15)
        tf.margin_top = Pt(15)
        
        # Populate checklist constraints items
        p_header = tf.paragraphs[0]
        p_header.text = "DESIGN CONSTRAINTS ANALYSIS"
        p_header.runs[0].font.name = "Arial"
        p_header.runs[0].font.size = Pt(11)
        p_header.runs[0].font.bold = True
        p_header.runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#37474F")
        
        for idx, constraint in enumerate(design_constraints_list, start=1):
            p = tf.add_paragraph()
            p.text = f"[✔] {constraint}"
            p.space_before = Pt(6)
            for run in p.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(10)
                run.font.color.rgb = ColorTranslator.hex_to_rgb("#455A64")
                
        logger.info("Engineering constraints analysis block generated.")
        return bg_card