import logging
from pptx import Presentation
from ..core_canvas.canvas_orchestrator import CanvasOrchestrator
from ..core_canvas.master_layout_mapper import MasterLayoutMapper
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("defense-presentation-expert")

class DefensePresentationExpert:
    def __init__(self, target_ratio="16_9"):
        self.canvas = CanvasOrchestrator(ratio=target_ratio)
        self.layout_mapper = MasterLayoutMapper()

    def compile_expert_defense_deck(self, fyp_project_title, candidate_name, supervisor_name, sections_manifest_list):
        logger.info(f"Compiling academic thesis defense presentation deck: {fyp_project_title}")
        
        # 1. Compile primary Title slide elements
        title_slide_layout = self.canvas.prs.slide_layouts[0] # Index 0 represents title slide layout
        title_slide = self.canvas.prs.slides.add_slide(title_slide_layout)
        
        self.layout_mapper.populate_placeholder_text(title_slide, "TITLE", fyp_project_title.upper())
        self.layout_mapper.populate_placeholder_text(title_slide, "SUBTITLE", fyp_project_title.upper())
        
        # Append candidate and supervisor details inside subtitle placeholder cell
        credentials_text = f"Candidate: {candidate_name}\nSupervisor: {supervisor_name}"
        self.layout_mapper.populate_placeholder_text(title_slide, "SUBTITLE", credentials_text)

        # 2. Compile secondary dynamic Section Divider slide frames
        for idx, section in enumerate(sections_manifest_list, start=1):
            divider_layout = self.canvas.prs.slide_layouts[2] # Divider layout index
            divider_slide = self.canvas.prs.slides.add_slide(divider_layout)
            
            self.layout_mapper.populate_placeholder_text(divider_slide, "TITLE", f"SECTION {idx:02D}")
            self.layout_mapper.populate_placeholder_text(divider_slide, "BODY", section.get("title", "").upper())

        logger.info("Expert academic thesis defense presentation deck compiled successfully.")
        return self.canvas