import logging
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("financial-charts")

class FinancialCharts:
    @staticmethod
    def draw_custom_gantt_chart(slide, left_emu, top_emu, width_emu, height_emu, tasks_list):
        logger.info(f"Compiling visual Gantt timeline chart. Tasks: {len(tasks_list)}")
        shapes = slide.shapes
        
        # Draw background timeline track card
        bg_card = shapes.add_shape(MSO_SHAPE.RECTANGLE, left_emu, top_emu, width_emu, height_emu)
        bg_card.fill.solid()
        bg_card.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#FAFAFA")
        bg_card.line.color.rgb = ColorTranslator.hex_to_rgb("#E0E0E0")
        
        step_height = int(height_emu / (len(tasks_list) + 1))
        
        # Render Gantt timeline task bars
        for idx, task in enumerate(tasks_list):
            task_name = task.get("name", f"Task {idx}")
            start_percent = task.get("start_percent", 0.1) # 0.0 to 1.0 range start
            duration_percent = task.get("duration_percent", 0.3)
            
            bar_left = left_emu + int(width_emu * start_percent)
            bar_width = int(width_emu * duration_percent)
            bar_top = top_emu + (idx + 1) * step_height - int(step_height * 0.25)
            bar_height = int(step_height * 0.5)
            
            # Create a rounded rectangular shape representing the task duration bar
            task_bar = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, bar_top, bar_width, bar_height)
            task_bar.fill.solid()
            task_bar.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#003366")
            task_bar.line.fill.background()
            
            task_bar.text = task_name
            for p in task_bar.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = ColorTranslator.hex_to_rgb("#FFFFFF")
                    
        logger.info("Custom Gantt timeline generated successfully.")
        return bg_card