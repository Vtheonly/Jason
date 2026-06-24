import logging
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("academic-analytics")

class AcademicAnalytics:
    @staticmethod
    def draw_confusion_matrix(slide, left_emu, top_emu, width_emu, height_emu, matrix_data):
        logger.info("Rendering scientific confusion matrix heatmap grid.")
        shapes = slide.shapes
        
        # Expected dimension layout: 2x2 grid calculations
        grid_dim = 2
        block_width = int(width_emu / grid_dim)
        block_height = int(height_emu / grid_dim)
        
        for r in range(grid_dim):
            for c in range(grid_dim):
                value = matrix_data[r][c]
                
                cell_left = left_emu + c * block_width
                cell_top = top_emu + r * block_height
                
                # Saturated heatmap color intensity calculations
                intensity_hex = "#E3F2FD" if value < 50 else "#1565C0"
                text_color_hex = "#0D47A1" if value < 50 else "#FFFFFF"
                
                cell_block = shapes.add_shape(MSO_SHAPE.RECTANGLE, cell_left, cell_top, block_width, block_height)
                cell_block.fill.solid()
                cell_block.fill.fore_color.rgb = ColorTranslator.hex_to_rgb(intensity_hex)
                cell_block.line.color.rgb = ColorTranslator.hex_to_rgb("#90CAF9")
                
                cell_block.text = str(value)
                for p in cell_block.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.name = "Arial"
                        run.font.size = Pt(16)
                        run.font.bold = True
                        run.font.color.rgb = ColorTranslator.hex_to_rgb(text_color_hex)
                        
        logger.info("Confusion matrix heatmap rendered.")