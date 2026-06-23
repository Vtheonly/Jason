import logging
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

logger = logging.getLogger("list-formatter")

class ListFormatter:
    @staticmethod
    def format_multi_level_bullets(text_frame, items_list, font_name="Calibri", base_size_pt=14, bullet_char="•"):
        logger.info(f"Formatting text frame lists sequence. Elements count: {len(items_list)}")
        try:
            # Clear default paragraphs template
            text_frame.clear()
            
            for idx, item in enumerate(items_list):
                # Format level depths based on leading dashes or tabulation spaces
                level = 0
                clean_text = item
                if item.startswith("  - "):
                    level = 1
                    clean_text = item[4:]
                elif item.startswith("    - "):
                    level = 2
                    clean_text = item[6:]
                elif item.startswith("- "):
                    clean_text = item[2:]

                p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
                p.text = clean_text
                p.level = level
                
                # Apply spacing and indent configurations based on nesting level
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                
                # Sizing hierarchy decay per list level
                target_size = base_size_pt - (level * 2)
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = Pt(max(target_size, 10))
                    
            logger.info("List formatting completed.")
        except Exception as err:
            logger.error(f"List rendering format broke: {str(err)}")