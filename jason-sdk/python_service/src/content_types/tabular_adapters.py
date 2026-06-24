import logging
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("tabular-adapters")

class TabularAdapters:
    @staticmethod
    def render_json_table(slide, left_emu, top_emu, width_emu, height_emu, header_list, rows_data):
        logger.info(f"Rendering table grid structure on slide. Columns: {len(header_list)}, Rows: {len(rows_data)}")
        try:
            shapes = slide.shapes
            rows_count = len(rows_data) + 1
            cols_count = len(header_list)
            
            # Create native styled table shape
            table_shape = shapes.add_table(rows_count, cols_count, left_emu, top_emu, width_emu, height_emu)
            table = table_shape.table
            
            # 1. Format and populate header row elements
            for c_idx, header_title in enumerate(header_list):
                cell = table.cell(0, c_idx)
                cell.text = str(header_title)
                cell.fill.solid()
                cell.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#003366")
                
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = ColorTranslator.hex_to_rgb("#FFFFFF")

            # 2. Format and populate alternating data rows elements
            for r_idx, row_values in enumerate(rows_data, start=1):
                bg_band = "#F9F9F9" if r_idx % 2 == 0 else "#FFFFFF"
                
                for c_idx, val in enumerate(row_values):
                    if c_idx >= cols_count: break
                    
                    cell = table.cell(r_idx, c_idx)
                    cell.text = str(val)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ColorTranslator.hex_to_rgb(bg_band)
                    
                    for run in cell.text_frame.paragraphs[0].runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(10)
                        run.font.color.rgb = ColorTranslator.hex_to_rgb("#333333")
                        
            logger.info("Table components generated successfully.")
            return table_shape
        except Exception as err:
            logger.error(f"Dynamic table matrix compilation failed: {str(err)}")
            raise err