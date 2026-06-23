import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("callout-warning-blocks")

class CalloutWarningBlocks:
    @staticmethod
    def draw_notice_card(slide, left_emu, top_emu, width_emu, height_emu, notice_text, level="info"):
        logger.info(f"Rendering stylized notice callout block. Severity Level: {level}")
        
        # Severity configuration profiles
        severity_theme = {
            "info": {"bg": "#EBF3FB", "border": "#2196F3", "text_color": "#0D47A1"},
            "warning": {"bg": "#FFF8E1", "border": "#FFC107", "text_color": "#E65100"},
            "danger": {"bg": "#FFEBEE", "border": "#F44336", "text_color": "#B71C1C"}
        }
        
        active = severity_theme.get(level, severity_theme["info"])
        
        shapes = slide.shapes
        bg_card = shapes.add_shape(MSO_SHAPE.RECTANGLE, left_emu, top_emu, width_emu, height_emu)
        
        # Style layout colors
        bg_card.fill.solid()
        bg_card.fill.fore_color.rgb = ColorTranslator.hex_to_rgb(active["bg"])
        bg_card.line.color.rgb = ColorTranslator.hex_to_rgb(active["border"])
        bg_card.line.width = Pt(2.0)
        
        # Populate notice text runs
        tf = bg_card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        
        p = tf.paragraphs[0]
        p.text = f"[{level.upper()}] {notice_text}"
        
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = ColorTranslator.hex_to_rgb(active["text_color"])
            
        logger.info("Notice callout card rendered successfully.")
        return bg_card