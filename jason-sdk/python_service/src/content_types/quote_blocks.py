import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("quote-blocks")

class QuoteBlocks:
    @staticmethod
    def draw_quote_card(slide, left_emu, top_emu, width_emu, height_emu, quote_body, author_name=None):
        logger.info("Generating stylized quote container frame.")
        shapes = slide.shapes
        
        # 1. Render quote block background
        quote_card = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_emu, top_emu, width_emu, height_emu)
        quote_card.fill.solid()
        quote_card.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#F5F5F5")
        quote_card.line.color.rgb = ColorTranslator.hex_to_rgb("#E0E0E0")
        quote_card.line.width = Pt(1.0)
        
        # Sane corner radius setting
        if hasattr(quote_card, 'adjustments') and len(quote_card.adjustments) > 0:
            quote_card.adjustments[0] = 0.04
            
        # 2. Draw stylized background quote graphic decoration
        quote_char = shapes.add_textbox(left_emu + Inches(0.15), top_emu - Inches(0.2), Inches(1), Inches(1))
        p_char = quote_char.text_frame.paragraphs[0]
        p_char.text = "“"
        p_char.runs[0].font.size = Pt(48)
        p_char.runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#BDBDBD")
        
        # 3. Populate quote text runs
        tf = quote_card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.4)
        tf.margin_right = Inches(0.4)
        tf.margin_top = Inches(0.3)
        tf.margin_bottom = Inches(0.3)
        
        p = tf.paragraphs[0]
        p.text = f"\"{quote_body}\""
        for run in p.runs:
            run.font.name = "Georgia"
            run.font.size = Pt(13)
            run.font.italic = True
            run.font.color.rgb = ColorTranslator.hex_to_rgb("#212121")
            
        # Append author name tag if present
        if author_name:
            p_author = tf.add_paragraph()
            p_author.text = f"— {author_name}"
            p_author.space_before = Pt(8)
            for run in p_author.runs:
                run.font.name = "Arial"
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.color.rgb = ColorTranslator.hex_to_rgb("#616161")
                
        logger.info("Quote container card generated successfully.")
        return quote_card