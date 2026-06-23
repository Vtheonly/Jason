import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("kpi-metric-cards")

class KpiMetricCards:
    @staticmethod
    def draw_kpi_card(slide, left_emu, top_emu, width_emu, height_emu, numeric_value, label_text, subtext=None):
        logger.info(f"Rendering KPI metric numeric display dashboard. Metric: {numeric_value}")
        
        shapes = slide.shapes
        card_shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_emu, top_emu, width_emu, height_emu)
        
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#FAFAFA")
        card_shape.line.color.rgb = ColorTranslator.hex_to_rgb("#EEEEEE")
        card_shape.line.width = Pt(1.5)
        
        # Sane corner radius setting
        if hasattr(card_shape, 'adjustments') and len(card_shape.adjustments) > 0:
            card_shape.adjustments[0] = 0.05
            
        tf = card_shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        
        # 1. Main Large KPI Numeric runs
        p_val = tf.paragraphs[0]
        p_val.text = str(numeric_value)
        p_val.runs[0].font.name = "Arial Black"
        p_val.runs[0].font.size = Pt(36)
        p_val.runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#003366")
        
        # 2. Main structural labeling
        p_label = tf.add_paragraph()
        p_label.text = str(label_text).upper()
        p_label.space_before = Pt(4)
        for run in p_label.runs:
            run.font.name = "Arial"
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = ColorTranslator.hex_to_rgb("#757575")
            
        # 3. Supplemental metric trend subtexts
        if subtext:
            p_sub = tf.add_paragraph()
            p_sub.text = str(subtext)
            p_sub.space_before = Pt(4)
            for run in p_sub.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(9)
                run.font.color.rgb = ColorTranslator.hex_to_rgb("#9E9E9E")
                
        logger.info("KPI metric card generated successfully.")
        return card_shape