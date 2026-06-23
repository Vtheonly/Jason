import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("check-notes-badges")

class CheckNotesBadges:
    @staticmethod
    def draw_status_badge(slide, left_emu, top_emu, width_emu, height_emu, badge_text, style="success"):
        logger.info(f"Rendering operational status tag. Value: {badge_text}, Style: {style}")
        
        # Visual styling config profiles
        style_config = {
            "success": {"bg": "#E8F5E9", "text": "#2E7D32"},
            "warning": {"bg": "#FFF3E0", "text": "#EF6C00"},
            "neutral": {"bg": "#ECEFF1", "text": "#37474F"}
        }
        
        active = style_config.get(style, style_config["neutral"])
        
        shapes = slide.shapes
        badge_shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_emu, top_emu, width_emu, height_emu)
        
        badge_shape.fill.solid()
        badge_shape.fill.fore_color.rgb = ColorTranslator.hex_to_rgb(active["bg"])
        badge_shape.line.fill.background()
        
        # Force curved stadium/pill shape corners
        if hasattr(badge_shape, 'adjustments') and len(badge_shape.adjustments) > 0:
            badge_shape.adjustments[0] = 0.5
            
        tf = badge_shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(badge_text).upper()
        
        for run in p.runs:
            run.font.name = "Arial"
            run.font.size = Pt(8)
            run.font.bold = True
            run.font.color.rgb = ColorTranslator.hex_to_rgb(active["text"])
            
        return badge_shape