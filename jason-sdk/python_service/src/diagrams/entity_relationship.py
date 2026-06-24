import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("entity-relationship")

class EntityRelationship:
    @staticmethod
    def draw_entity_relationship_diagram(slide, left_emu, top_emu, width_emu, height_emu, tables):
        logger.info(f"Generating Entity Relationship database diagrams tables. Tables: {len(tables)}")
        shapes = slide.shapes
        
        tables_count = len(tables)
        if tables_count == 0: return

        block_width = int(width_emu / tables_count) - Pt(15)
        
        for idx, table in enumerate(tables):
            table_name = table.get("name", "Table")
            columns = table.get("columns", [])

            col_left = left_emu + idx * (block_width + Pt(15))
            
            # Draw table structure rectangle
            table_box = shapes.add_shape(MSO_SHAPE.RECTANGLE, col_left, top_emu, block_width, height_emu)
            table_box.fill.solid()
            table_box.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#FAFAFA")
            table_box.line.color.rgb = ColorTranslator.hex_to_rgb("#003366")
            table_box.line.width = Pt(1.5)

            # Format and populate column rows
            tf = table_box.text_frame
            tf.clear()
            
            p_header = tf.paragraphs[0]
            p_header.text = table_name.upper()
            p_header.runs[0].font.bold = True
            p_header.runs[0].font.size = Pt(10)
            p_header.runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#003366")

            p_sep = tf.add_paragraph()
            p_sep.text = "===================="

            for col in columns:
                col_name = col.get("name", "col")
                col_type = col.get("type", "VARCHAR")
                is_pk = col.get("primary_key", False)

                p_col = tf.add_paragraph()
                pk_indicator = "[PK] " if is_pk else ""
                p_col.text = f"{pk_indicator}{col_name} : {col_type}"
                
                # Format text runs style options
                for r in p_col.runs:
                    r.font.name = "Consolas"
                    r.font.size = Pt(8)
                    r.font.color.rgb = ColorTranslator.hex_to_rgb("#333333")
                    if is_pk: r.font.bold = True
                    
        logger.info("Entity relationship diagram successfully updated.")