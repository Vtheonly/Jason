import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("pipeline-visualizers")

class PipelineVisualizers:
    @staticmethod
    def draw_data_pipeline_flow(slide, left_emu, top_emu, width_emu, height_emu, pipeline_steps):
        logger.info(f"Rendering visual timeline for data pipeline flow steps: {len(pipeline_steps)}")
        shapes = slide.shapes
        
        steps_count = len(pipeline_steps)
        if steps_count == 0: return

        block_width = int(width_emu / steps_count) - Pt(10)
        block_height = int(height_emu * 0.5)

        for idx, step_name in enumerate(pipeline_steps):
            step_left = left_emu + idx * (block_width + Pt(10))
            step_top = top_emu + int((height_emu - block_height) / 2)

            # Standard chevron step indicator box
            step_box = shapes.add_shape(MSO_SHAPE.CHEVRON, step_left, step_top, block_width, block_height)
            step_box.fill.solid()
            step_box.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#4682B4")
            step_box.line.color.rgb = ColorTranslator.hex_to_rgb("#003366")

            # Format and populate text blocks
            step_box.text = step_name
            for p in step_box.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(10)
                    run.font.bold = True
                    run.font.color.rgb = ColorTranslator.hex_to_rgb("#FFFFFF")

        logger.info("Pipeline visual workflow blocks compiled successfully.")