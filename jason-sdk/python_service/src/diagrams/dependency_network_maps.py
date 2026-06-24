import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("dependency-network-maps")

class DependencyNetworkMaps:
    @staticmethod
    def draw_layered_network(slide, left_emu, top_emu, width_emu, height_emu, nodes, edges):
        logger.info(f"Rendering system dependency network layers map. Nodes count: {len(nodes)}")
        shapes = slide.shapes
        
        node_positions = {}
        nodes_count = len(nodes)
        if nodes_count == 0: return

        # Multi-layer circular placements equations
        import math
        center_x = left_emu + int(width_emu / 2)
        center_y = top_emu + int(height_emu / 2)
        radius = int(min(width_emu, height_emu) * 0.35)

        # 1. Render Circular Node Blocks
        for idx, node in enumerate(nodes):
            node_id = node.get("id")
            label = node.get("label", "Node")

            angle = (2 * math.pi * idx) / nodes_count
            node_left = center_x + int(radius * math.cos(angle)) - Pt(15)
            node_top = center_y + int(radius * math.sin(angle)) - Pt(15)

            box = shapes.add_shape(MSO_SHAPE.OVAL, node_left, node_top, Pt(30), Pt(30))
            box.fill.solid()
            box.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#003366")
            box.line.color.rgb = ColorTranslator.hex_to_rgb("#FFBF00")
            
            box.text = label
            for p in box.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(8)
                    run.font.bold = True
                    run.font.color.rgb = ColorTranslator.hex_to_rgb("#FFFFFF")

            node_positions[node_id] = (node_left + Pt(15), node_top + Pt(15))

        # 2. Render connecting lines mapping relationships
        for edge in edges:
            from_node = edge.get("from")
            to_node = edge.get("to")

            if from_node in node_positions and to_node in node_positions:
                start_pt = node_positions[from_node]
                end_pt = node_positions[to_node]

                # Draw thin connection lines between coordinates
                line_width = int(end_pt[0] - start_pt[0])
                line_height = int(end_pt[1] - start_pt[1])
                
                connector = shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 
                    int(start_pt[0]), int(start_pt[1]), 
                    abs(line_width) or Pt(1), abs(line_height) or Pt(1)
                )
                connector.fill.solid()
                connector.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#B0BEC5")
                connector.line.fill.background()
                
        logger.info("Layered network map generated.")