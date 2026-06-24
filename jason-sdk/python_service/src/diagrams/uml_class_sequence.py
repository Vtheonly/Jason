import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("uml-class-sequence")

class UmlClassSequence:
    @staticmethod
    def draw_class_blocks(slide, left_emu, top_emu, width_emu, height_emu, classes):
        logger.info("Generating standard UML structural class diagrams partitions.")
        shapes = slide.shapes
        
        class_count = len(classes)
        if class_count == 0: return

        block_width = int(width_emu / class_count) - Pt(10)
        
        for idx, cls in enumerate(classes):
            name = cls.get("label", "Class")
            attributes = cls.get("attributes", [])
            methods = cls.get("methods", [])

            col_left = left_emu + idx * (block_width + Pt(10))
            
            # Construct distinct partitions rectangles
            class_card = shapes.add_shape(MSO_SHAPE.RECTANGLE, col_left, top_emu, block_width, height_emu)
            class_card.fill.solid()
            class_card.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#F5F5F5")
            class_card.line.color.rgb = ColorTranslator.hex_to_rgb("#003366")
            class_card.line.width = Pt(1.5)

            # Format and populate text blocks
            tf = class_card.text_frame
            tf.clear()
            
            # Header Row
            p_name = tf.paragraphs[0]
            p_name.text = f"<<{name}>>"
            p_name.alignment = 1 # Horizontal centering
            p_name.runs[0].font.bold = True
            p_name.runs[0].font.size = Pt(11)
            p_name.runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#003366")

            # Attributes block
            p_sep1 = tf.add_paragraph()
            p_sep1.text = "------------------------"
            for attr in attributes:
                p_attr = tf.add_paragraph()
                p_attr.text = f"- {attr}"
                p_attr.runs[0].font.size = Pt(9)
                p_attr.runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#333333")

            # Methods block
            p_sep2 = tf.add_paragraph()
            p_sep2.text = "------------------------"
            for meth in methods:
                p_meth = tf.add_paragraph()
                p_meth.text = f"+ {meth}()"
                p_meth.runs[0].font.size = Pt(9)
                p_meth.runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#333333")

        logger.info("UML Class diagram generated successfully.")

    @staticmethod
    def draw_sequence_diagram(slide, left_emu, top_emu, width_emu, height_emu, actors, messages):
        logger.info("Generating standard sequence diagram lifeline tracks.")
        shapes = slide.shapes
        
        # 1. Draw lifelines vertical tracks
        actors_count = len(actors)
        if actors_count == 0: return

        spacing = int(width_emu / actors_count)
        lifeline_map = {}

        for idx, act in enumerate(actors):
            act_id = act.get("id")
            label = act.get("label", "Actor")

            center_x = left_emu + idx * spacing + int(spacing / 2)
            lifeline_map[act_id] = center_x

            # Lifeline header box shape
            header_box = shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Pt(30), top_emu, Pt(60), Pt(20))
            header_box.fill.solid()
            header_box.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#4682B4")
            header_box.text = label
            header_box.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
            header_box.text_frame.paragraphs[0].runs[0].font.color.rgb = ColorTranslator.hex_to_rgb("#FFFFFF")

            # Lifeline vertical guide dotted line
            shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Pt(1), top_emu + Pt(20), Pt(2), height_emu - Pt(20))

        # 2. Draw message passing connectors
        for idx, msg in enumerate(messages):
            from_act = msg.get("from")
            to_act = msg.get("to")
            text = msg.get("label", "")

            if from_act in lifeline_map and to_act in lifeline_map:
                start_x = lifeline_map[from_act]
                end_x = lifeline_map[to_act]
                
                step_y = top_emu + Pt(40) + idx * Pt(30)
                arrow_width = end_x - start_x
                
                # Dynamic horizontal messaging arrows
                arrow = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, start_x, step_y, arrow_width, Pt(3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#FFBF00")
                arrow.line.fill.background()