import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from ..theme_engine.color_translator import ColorTranslator

logger = logging.getLogger("flowchart-blocks")

class FlowchartBlocks:
    @staticmethod
    def draw_flowchart(slide, left_emu, top_emu, width_emu, height_emu, nodes, edges):
        logger.info(f"Compiling vector flowchart block. Nodes count: {len(nodes)}")
        shapes = slide.shapes
        
        node_map = {}
        nodes_count = len(nodes)
        if nodes_count == 0: return

        # Simple horizontal/vertical layout placement matrices calculations
        block_width = int(width_emu / (nodes_count + 1))
        block_height = int(height_emu * 0.4)
        
        # 1. Render flowchart process boxes
        for idx, node in enumerate(nodes):
            node_id = node.get("id")
            label = node.get("label", "")
            node_type = node.get("type", "process")

            shape_type = MSO_SHAPE.RECTANGLE
            if node_type == "start" or node_type == "end":
                shape_type = MSO_SHAPE.OVAL
            elif node_type == "decision":
                shape_type = MSO_SHAPE.DIAMOND

            node_left = left_emu + (idx * block_width) + int(block_width * 0.1)
            node_top = top_emu + int((height_emu - block_height) / 2)
            
            box = shapes.add_shape(shape_type, node_left, node_top, int(block_width * 0.8), block_height)
            box.fill.solid()
            box.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#003366")
            box.line.color.rgb = ColorTranslator.hex_to_rgb("#4682B4")
            
            box.text = label
            for p in box.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = ColorTranslator.hex_to_rgb("#FFFFFF")

            node_map[node_id] = {
                "left": node_left,
                "top": node_top,
                "width": int(block_width * 0.8),
                "height": block_height
            }

        # 2. Render connecting lines between steps
        for edge in edges:
            from_node = edge.get("from")
            to_node = edge.get("to")
            if from_node in node_map and to_node in node_map:
                # Draw direct horizontal connector lines
                start = node_map[from_node]
                end = node_map[to_node]
                
                line_left = start["left"] + start["width"]
                line_width = end["left"] - line_left
                line_top = start["top"] + int(start["height"] / 2)
                
                connector = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, line_left, line_top - Pt(2), line_width, Pt(4))
                connector.fill.solid()
                connector.fill.fore_color.rgb = ColorTranslator.hex_to_rgb("#FFBF00")
                connector.line.fill.background()
                
        logger.info("Vector flowchart rendering completed.")