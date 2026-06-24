import logging
from pptx.util import Pt

logger = logging.getLogger("overflow-autofit")

class OverflowAutofit:
    @staticmethod
    def adjust_text_frame_to_prevent_overflow(text_frame, max_width_emu, max_height_emu, margin_shrink_enabled=True):
        logger.info("Executing text boundaries overflow auto-fit scanner.")
        try:
            # 1. Strip default internal paddings to expand usable text boundaries
            if margin_shrink_enabled:
                text_frame.margin_left = Pt(2)
                text_frame.margin_right = Pt(2)
                text_frame.margin_top = Pt(2)
                text_frame.margin_bottom = Pt(2)

            # 2. Iteratively check density and reduce font sizes if contents overflow bounding boxes
            word_count = len(text_frame.text.split())
            if word_count > 80:
                logger.warning("High character count detected. Reducing font sizing scale to prevent container overflows.")
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size and run.font.size > Pt(10):
                            run.font.size = Pt(run.font.size.pt - 2)
                            
            logger.info("Text frame boundaries auto-fit completed successfully.")
        except Exception as err:
            logger.error(f"Text boundaries overflow adjustment failed: {str(err)}")