import logging
from pptx import Presentation

logger = logging.getLogger("outline-extractor")

class OutlineExtractor:
    @staticmethod
    def extract_text_and_notes_outlines(pptx_file_path, output_text_path):
        logger.info(f"Extracting textual items outlines from compiled file: {pptx_file_path}")
        try:
            prs = Presentation(pptx_file_path)
            outline_data = []

            for idx, slide in enumerate(prs.slides):
                slide_outline = {
                    "slide_index": idx,
                    "title": "",
                    "bullets": [],
                    "notes": ""
                }

                # 1. Extract plain text bullet content items
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text.strip()
                        if not text: continue
                        
                        # Use simple heuristics to separate heading elements from list runs
                        if shape.name.lower().includes("title") or shape.height < 500000:
                            slide_outline["title"] = text
                        else:
                            slide_outline["bullets"].extend([line.strip() for line in text.split('\n') if line.strip()])

                # 2. Extract slide notes paragraphs
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    slide_outline["notes"] = slide.notes_slide.notes_text_frame.text.strip()

                outline_data.append(slide_outline)

            # Write text summaries directly to target text summary file
            with open(output_text_path, "w", encoding="utf-8") as f:
                f.write("================================================================\n")
                f.write("               PRESENTATION OUTLINE EXPORT SUMMARY\n")
                f.write("================================================================\n")
                
                for slide in outline_data:
                    f.write(f"\n[Slide Index: {slide['slide_index']}] | Title: \"{slide['title']}\"\n")
                    for b in slide["bullets"]:
                        f.write(f"  • {b}\n")
                    if slide["notes"]:
                        f.write(f"  Presenter Notes: {slide['notes']}\n")

            logger.info(f"Outline details written to summary target path: {output_text_path}")
            return outline_data
        except Exception as err:
            logger.error("Outline elements extraction failed.", exc_info=True)
            raise err