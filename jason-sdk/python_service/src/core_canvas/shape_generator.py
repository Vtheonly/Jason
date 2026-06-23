import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger("shape-generator")

class ShapeGenerator:
    @staticmethod
    def draw_rectangle(slide, left_emu, top_emu, width_emu, height_emu, fill_color=None, border_color=None):
        logger.info("Adding rectangular shape parameters to target slide layout.")
        shapes = slide.shapes
        shape = shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left_emu, top_emu, width_emu, height_emu
        )
        
        # Configure fill color parameters
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()

        # Configure border styling properties
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()

        return shape

    @staticmethod
    def draw_rounded_rectangle(slide, left_emu, top_emu, width_emu, height_emu, fill_color=None, border_color=None, radius=0.1):
        logger.info("Adding rounded rectangular container parameters.")
        shapes = slide.shapes
        shape = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left_emu, top_emu, width_emu, height_emu
        )
        
        # Adjust curve corners factor
        if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
            shape.adjustments[0] = radius

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)

        return shape