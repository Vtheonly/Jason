import os
import logging
from pptx import Presentation
from pptx.util import Inches
from core_canvas.coordinate_transformer import CoordinateTransformer

logger = logging.getLogger("canvas-orchestrator")

class CanvasOrchestrator:
    def __init__(self, pptx_path=None, ratio="16_9"):
        self.pptx_path = pptx_path
        self.ratio = ratio
        
        if pptx_path and os.path.exists(pptx_path):
            logger.info(f"Opening existing presentation canvas: {pptx_path}")
            self.prs = Presentation(pptx_path)
        else:
            logger.info("Initializing fresh presentation blank canvas.")
            self.prs = Presentation()
            self._set_canvas_dimensions()

    def _set_canvas_dimensions(self):
        # Configure standard aspect ratio bounds
        if self.ratio == "16_9":
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        else:
            self.prs.slide_width = Inches(10.0)
            self.prs.slide_height = Inches(7.5)
        logger.info(f"Canvas size initialized. Aspect: {self.ratio} (W:{self.prs.slide_width.inches}\", H:{self.prs.slide_height.inches}\")")

    def add_blank_slide(self):
        # Index 6 of default presentation themes is blank layout slide
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        logger.info(f"Appended blank slide. Total slides count: {len(self.prs.slides)}")
        return slide

    def save(self, target_output_path):
        logger.info(f"Saving compiled presentation layout at: {target_output_path}")
        self.prs.save(target_output_path)