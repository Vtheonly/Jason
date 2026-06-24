"""Agenda slide builder."""
import logging
from pptx.util import Pt
from content_types.list_formatter import ListFormatter
from core_canvas.coordinate_transformer import CoordinateTransformer

logger = logging.getLogger("agenda-slide-builder")


class AgendaSlideBuilder:
    @staticmethod
    def build(canvas, agenda_items):
        logger.info(f"Building agenda slide with {len(agenda_items)} items.")
        blank_layout = canvas.prs.slide_layouts[6]
        slide = canvas.prs.slides.add_slide(blank_layout)

        left = CoordinateTransformer.inches_to_emu(0.8)
        top = CoordinateTransformer.inches_to_emu(1.5)
        width = CoordinateTransformer.inches_to_emu(11.5)
        height = CoordinateTransformer.inches_to_emu(5.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        ListFormatter.format_multi_level_bullets(tf, agenda_items)
        return slide
