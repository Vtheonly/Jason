"""Title slide builder."""
import logging
from pptx.util import Pt, Inches
from theme_engine.color_translator import ColorTranslator
from theme_engine.font_pairer import FontPairer

logger = logging.getLogger("title-slide-builder")


class TitleSlideBuilder:
    @staticmethod
    def build(canvas, title_text, subtitle_text=None, theme=None):
        logger.info(f"Building title slide: {title_text}")
        slide_layout = canvas.prs.slide_layouts[0]
        slide = canvas.prs.slides.add_slide(slide_layout)

        for ph in slide.placeholders:
            ph_idx = ph.placeholder_format.idx
            if ph_idx == 0 and title_text:
                ph.text = title_text
                for run in ph.text_frame.paragraphs[0].runs:
                    FontPairer.apply_text_styling(run, "Arial", 36, is_bold=True)
            elif ph_idx == 1 and subtitle_text:
                ph.text = subtitle_text
                for run in ph.text_frame.paragraphs[0].runs:
                    FontPairer.apply_text_styling(run, "Calibri", 18, is_italic=True)

        return slide
