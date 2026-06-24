"""Closing & appendix slides builder."""
import logging
from pptx.util import Pt
from core_canvas.coordinate_transformer import CoordinateTransformer

logger = logging.getLogger("closing-appendix-slides")


class ClosingAppendixSlidesBuilder:
    @staticmethod
    def build(canvas, thanks_text="Thank You", appendix_items=None):
        logger.info("Building closing/appendix slide.")
        title_layout = canvas.prs.slide_layouts[0]
        slide = canvas.prs.slides.add_slide(title_layout)

        for ph in slide.placeholders:
            ph_idx = ph.placeholder_format.idx
            if ph_idx == 0:
                ph.text = thanks_text
            elif ph_idx == 1 and appendix_items:
                ph.text = "\n".join(appendix_items)
        return slide
